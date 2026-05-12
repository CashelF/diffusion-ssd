from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0xC8, 0x1D, 0x25)
GREY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xAA, 0xAA, 0xAA)


def add_title(slide, text, top=0.35):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    return box


def add_bullets(slide, bullets, top=1.6, left=0.7, width=12.0, height=5.5,
                size=22, color=GREY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.size = Pt(size if level == 0 else max(14, size - 4))
        p.font.color.rgb = color if level == 0 else GREY
        p.space_after = Pt(8)
    return box


def add_accent_bar(slide, top=1.25):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(top),
                                 Inches(1.2), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


blank = prs.slide_layouts[6]

# Slide 0: Title
s0 = prs.slides.add_slide(blank)
title_box = s0.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.6))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Diffusion Drafters for Speculative Squared Decoding"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = NAVY

bar = s0.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.25),
                          Inches(1.5), Inches(0.1))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

sub = s0.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.5))
tfs = sub.text_frame
tfs.word_wrap = True
ps = tfs.paragraphs[0]
ps.text = "Parallelizing the draft model to accelerate large language model inference"
ps.font.size = Pt(24)
ps.font.color.rgb = GREY

ps2 = tfs.add_paragraph()
ps2.text = "CS5788 Generative Models, Cornell"
ps2.font.size = Pt(18)
ps2.font.color.rgb = LIGHT
ps2.space_before = Pt(12)

# Slide 1: Problem
s1 = prs.slides.add_slide(blank)
add_title(s1, "Problem: LLM inference is bottlenecked by autoregressive decoding")
add_accent_bar(s1)
add_bullets(s1, [
    "Large LLMs generate one token at a time. Each token requires a full forward pass over billions of parameters.",
    "Speculative decoding (SD) speeds this up using two models: a small fast draft model proposes k candidate tokens, and the large target model verifies them.",
    "Why verification is fast: the target model runs all k draft tokens through one forward pass in parallel, comparing its own predicted distribution at each position to the drafter’s. Tokens are accepted up to the first mismatch.",
    "Result: instead of k sequential forward passes through the large model, you pay roughly one large pass per accepted run of tokens.",
    "But standard SD is still sequential overall. Drafting and verification alternate on the same hardware, so each is idle while the other runs.",
    "Speculative Squared Decoding (SSD) overlaps drafting and verification on separate hardware by pre‐speculating over the tree of likely verifier outcomes.",
    "Remaining bottleneck: the draft model is itself autoregressive, so generating each draft is a sequential, k‐step process.",
])

# Slide 2: Goal
s2 = prs.slides.add_slide(blank)
add_title(s2, "Goal: replace the autoregressive drafter with a diffusion drafter")
add_accent_bar(s2)
add_bullets(s2, [
    "Main task: SSD with an autoregressive verifier and a diffusion draft model.",
    "Diffusion language models (MDLM, BD3‐LM) generate all k draft tokens in parallel rather than left to right.",
    "Hypothesis: parallel drafting cuts draft latency, increases the speculation cache size SSD can maintain, and raises end‐to‐end throughput.",
    "Extension: SSD with a block‐diffusion verifier and a diffusion drafter, parallelizing both sides.",
    "Success metric: tokens per second and draft acceptance rate on humaneval, alpaca, gsm8k, and ultrafeedback compared to the AR‐drafter SSD baseline.",
])

# Slide 3: Approach
s3 = prs.slides.add_slide(blank)
add_title(s3, "Approach")
add_accent_bar(s3)
add_bullets(s3, [
    "Start from the SSD inference algorithm of Kumar, Dao, May (ICLR 2026) as the baseline.",
    "Swap the autoregressive draft model for a diffusion language model that produces all k draft tokens in a single parallel pass.",
    "Two diffusion sampling strategies:",
    ("Remask sampler with a configurable number of refinement steps.", 1),
    ("BD3‐LM staircase sampler with cached prefix reuse for streaming generation.", 1),
    "Optional autoregressive warm start: prefill the first N draft tokens with a small AR model, then hand off to the diffusion model.",
    "Evaluate against AR only, sync SD, async SSD with an AR drafter, and SGLang and vLLM speculative decoding servers.",
])

# Slide 4: Results
s4 = prs.slides.add_slide(blank)
add_title(s4, "Results: AR drafting wins at scale; BD3-LM drafter is the bottleneck")
add_accent_bar(s4)

# Subtitle / setup line
setup = s4.shapes.add_textbox(Inches(0.6), Inches(1.20), Inches(12.1), Inches(0.5))
tfx = setup.text_frame
tfx.word_wrap = True
px = tfx.paragraphs[0]
px.text = "1×H100, Qwen3 family, sync SD, temp=0, 512 output tokens; humaneval / alpaca / gsm8k / ultrafeedback."
px.font.size = Pt(14)
px.font.italic = True
px.font.color.rgb = LIGHT

# Build a results table
rows_data = [
    ("Target", "Drafter", "Throughput (tok/s)", "vs baseline"),
    ("Qwen3-1.7B", "AR baseline (no spec)", "393", "1.00×"),
    ("Qwen3-1.7B", "Sync SD + AR draft (Qwen3-0.6B, k=6)", "285", "0.72×"),
    ("Qwen3-1.7B", "Sync SD + BD3-LM diffusion draft (k=15)", "17", "0.04×"),
    ("Qwen3-8B", "AR baseline (no spec)", "147", "1.00×"),
    ("Qwen3-8B", "Sync SD + AR draft (Qwen3-0.6B, k=6)", "193", "1.31×"),
    ("Qwen3-8B", "Sync SD + BD3-LM diffusion draft (k=15)", "20", "0.13×"),
]
rows, cols = len(rows_data), len(rows_data[0])
table_shape = s4.shapes.add_table(rows, cols, Inches(0.6), Inches(1.75),
                                  Inches(9.4), Inches(2.8))
table = table_shape.table
table.columns[0].width = Inches(1.6)
table.columns[1].width = Inches(4.4)
table.columns[2].width = Inches(2.0)
table.columns[3].width = Inches(1.4)

for r, row_vals in enumerate(rows_data):
    for c, val in enumerate(row_vals):
        cell = table.cell(r, c)
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = val
        run = para.runs[0]
        run.font.size = Pt(13)
        if r == 0:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            run.font.color.rgb = GREY
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # highlight the only winning row
        if r == 5:
            run.font.bold = True
            run.font.color.rgb = ACCENT

# Findings panel on the right
findings = s4.shapes.add_textbox(Inches(10.2), Inches(1.75), Inches(2.95), Inches(5.0))
tff = findings.text_frame
tff.word_wrap = True
fp = tff.paragraphs[0]
fp.text = "Findings"
fp.font.size = Pt(16)
fp.font.bold = True
fp.font.color.rgb = NAVY
fp.space_after = Pt(6)

for line in [
    "AR-SD reaches its win regime at 8B target (1.31×). Confirms the setup is correct.",
    "At 1.7B target, even AR-SD is net-slower; draft overhead dominates.",
    "BD3-LM diffusion draft yields ~17–20 tok/s regardless of target size — the drafter itself is the bottleneck, not the verifier.",
    "Block backend lacks the CUDA-graph / JIT path the AR drafter uses; 4 refine steps × 0.6B params per cycle dominate latency.",
]:
    p = tff.add_paragraph()
    p.text = "•  " + line
    p.font.size = Pt(11)
    p.font.color.rgb = GREY
    p.space_after = Pt(6)

# Caveats footnote
foot = s4.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.7))
tfo = foot.text_frame
tfo.word_wrap = True
po = tfo.paragraphs[0]
po.text = (
    "Caveats: numseqs differs across runs (A/B = 512 prompts, A8/B8 = 256, C = 128, C8 = 64) — throughput is per-token average. "
    "Async SSD with diffusion drafter not measured: codebase requires sync mode for --draft-backend block."
)
po.font.size = Pt(10)
po.font.italic = True
po.font.color.rgb = LIGHT

out = "/Users/tishya/Desktop/Cornell/Y1S2/gen_models/project/code/diffusion-ssd/diffusion_ssd_slides.pptx"
prs.save(out)
print(out)
